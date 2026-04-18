"""
PowerPoint ingestor for Ariadne.

Extracts text from .pptx files using python-pptx.
"""

from pathlib import Path
from typing import List

from ariadne.ingest.base import BaseIngestor, Document, SourceType


class PPTIngestor(BaseIngestor):
    """
    Ingest PowerPoint (.pptx) presentations.

    Treats each slide as a chunk, concatenating title + body text.
    This preserves the slide's narrative structure.
    """

    source_type = SourceType.PPT

    def _extract(self, path: Path) -> List[str]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required to ingest .pptx files. "
                "Install it with: pip install python-pptx"
            )

        prs = Presentation(str(path))
        chunks = []

        for slide_num, slide in enumerate(prs.slides, 1):
            parts = []

            # Extract title
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                if title:
                    parts.append(f"Slide {slide_num} Title: {title}")

            # Extract body text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)

            if parts:
                chunks.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

        return chunks if chunks else []
